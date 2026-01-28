"""
Motor de generación de presentaciones PPTX para ProfeIC.
Convierte estructuras JSON validadas (Pydantic) en archivos .pptx.
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.prompts import PresentationStructure


def create_pptx_from_json(presentation_data: PresentationStructure) -> BytesIO:
    """
    Crea un archivo PPTX a partir de una estructura PresentationStructure validada.
    
    Args:
        presentation_data: Objeto PresentationStructure con los datos de la presentación
    
    Returns:
        BytesIO con el contenido del archivo .pptx
    """
    # Crear presentación
    prs = Presentation()
    
    # Configurar tamaño de diapositiva (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Crear diapositiva de título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = presentation_data.title
    subtitle.text = "Generado por ProfeIC"
    
    # Agregar cada diapositiva de contenido
    for slide_data in presentation_data.slides:
        # Usar layout de título y contenido
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Título de la diapositiva
        title_shape = slide.shapes.title
        title_shape.text = slide_data.slide_title
        
        # Contenido con viñetas
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        # Agregar puntos clave
        for i, bullet_point in enumerate(slide_data.bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = bullet_point
                p.level = 0
                p.font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = bullet_point
                p.level = 0
                p.font.size = Pt(18)
        
        # Agregar notas del orador
        notes_slide = slide.notes_slide
        notes_shape = notes_slide.notes_text_frame
        notes_shape.text = slide_data.speaker_notes
    
    # Guardar en BytesIO
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer


def create_pptx_from_dict(presentation_dict: dict) -> BytesIO:
    """
    Crea un archivo PPTX a partir de un diccionario (útil para datos de API).
    
    Args:
        presentation_dict: Diccionario con estructura de PresentationStructure
    
    Returns:
        BytesIO con el contenido del archivo .pptx
    """
    # Validar y convertir a Pydantic
    presentation_data = PresentationStructure(**presentation_dict)
    return create_pptx_from_json(presentation_data)

