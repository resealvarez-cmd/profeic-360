import io
import os
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptxRGB

# --- CONFIGURACIÓN ---
TEMPLATE_PPT = "formato ppt ia.pptx"
LOGO_PATH = "logo.png"
BLUE = RGBColor(26, 46, 59)

def add_header(doc, title):
    """Encabezado estándar"""
    section = doc.sections[0]
    header = section.header
    p = header.paragraphs[0]
    p.text = "COLEGIO MADRE PAULINA\n" + title
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

# --- 1. PLANIFICACIÓN ---
def create_plan_docx(data):
    doc = Document()
    add_header(doc, "PLANIFICACIÓN CURRICULAR")
    doc.add_heading(data.get('titulo_unidad', 'Unidad'), 0)
    
    for clase in data.get('clases', []):
        doc.add_heading(f"Clase {clase.get('numero')}", level=1)
        doc.add_paragraph(f"Meta: {clase.get('meta_clase','')}")
        doc.add_paragraph("Modelamiento:", style='Heading 3')
        doc.add_paragraph(clase.get('paso_3_modelamiento',''))
        doc.add_paragraph("Práctica Deliberada:", style='Heading 3')
        for t in clase.get('paso_5_practica_deliberada', []):
            doc.add_paragraph(f"• {t}")
        doc.add_page_break()
    
    s = io.BytesIO(); doc.save(s); s.seek(0); return s

def create_class_ppt(data):
    prs = Presentation(TEMPLATE_PPT) if os.path.exists(TEMPLATE_PPT) else Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get('clase', {}).get('meta_clase', 'Clase')
    s = io.BytesIO(); prs.save(s); s.seek(0); return s

# --- 2. RÚBRICA ---
def create_rubric_docx(data):
    doc = Document()
    add_header(doc, "RÚBRICA DE EVALUACIÓN")
    doc.add_heading(data.get('titulo', 'Rúbrica'), 0)
    doc.add_paragraph(data.get('descripcion', ''))
    
    table = doc.add_table(rows=1, cols=5)
    table.style = 'Table Grid'
    hdr = table.rows[0].cells
    hdr[0].text = "Criterio"
    hdr[1].text = "Insuficiente"
    hdr[2].text = "Elemental"
    hdr[3].text = "Adecuado"
    hdr[4].text = "Destacado"
    
    for crit in data.get('tabla', []):
        row = table.add_row().cells
        row[0].text = f"{crit.get('criterio')}\n({crit.get('porcentaje')}%)"
        niv = crit.get('niveles', {})
        row[1].text = niv.get('insuficiente','')
        row[2].text = niv.get('elemental','')
        row[3].text = niv.get('adecuado','')
        row[4].text = niv.get('destacado','')
        
    s = io.BytesIO(); doc.save(s); s.seek(0); return s

# --- 3. EVALUACIÓN ---
def create_assessment_docx(data):
    doc = Document()
    add_header(doc, "EVALUACIÓN ESCRITA")
    doc.add_heading(data.get('titulo', 'Prueba'), 0)
    doc.add_paragraph(data.get('description', ''))
    
    for i, item in enumerate(data.get('items', [])):
        p = doc.add_paragraph()
        p.add_run(f"{i+1}. {item.get('stem')}").bold = True
        p.add_run(f" ({item.get('points')} pts)")
        
        if item.get('options'):
            for opt in item['options']:
                # Handle both string and dict (legacy)
                text = opt.get('text') if isinstance(opt, dict) else opt
                doc.add_paragraph(f"   ○ {text}")
        else:
            doc.add_paragraph("\n" * 3) # Espacio para respuesta
            
    s = io.BytesIO(); doc.save(s); s.seek(0); return s
